"""
MS Office Skill - Read and edit Microsoft Office files.

Supports:
- Word (.docx): Read, create, edit documents
- Excel (.xlsx): Read, create, edit spreadsheets
- PowerPoint (.pptx): Read, create, edit presentations
"""

import json
import os
from datetime import datetime
from pathlib import Path
from typing import Optional

from r_cli.core.agent import Skill
from r_cli.core.config import Config
from r_cli.core.llm import Tool

# Optional imports
try:
    from docx import Document
    from docx.enum.text import WD_ALIGN_PARAGRAPH
    from docx.shared import Inches, Pt

    HAS_DOCX = True
except ImportError:
    HAS_DOCX = False

try:
    from openpyxl import Workbook, load_workbook
    from openpyxl.styles import Alignment, Border, Font, PatternFill, Side
    from openpyxl.utils import get_column_letter

    HAS_XLSX = True
except ImportError:
    HAS_XLSX = False

try:
    from pptx import Presentation
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Inches as PptxInches
    from pptx.util import Pt as PptxPt

    HAS_PPTX = True
except ImportError:
    HAS_PPTX = False


class MSOfficeSkill(Skill):
    """
    Microsoft Office file manipulation skill.

    Read, create, and edit Word, Excel, and PowerPoint files.
    """

    name = "msoffice"
    description = "Read and edit Microsoft Office files (Word, Excel, PowerPoint)"

    def __init__(self, config: Optional[Config] = None):
        super().__init__(config)

    def get_tools(self) -> list[Tool]:
        tools = []

        # Word tools
        if HAS_DOCX:
            tools.extend(
                [
                    Tool(
                        name="word_read",
                        description="Read content from a Word document (.docx)",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .docx file",
                                },
                                "include_tables": {
                                    "type": "boolean",
                                    "description": "Include tables in output (default: true)",
                                },
                            },
                            "required": ["file_path"],
                        },
                        handler=self.word_read,
                    ),
                    Tool(
                        name="word_create",
                        description="Create a new Word document with content",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Output path for the .docx file",
                                },
                                "title": {
                                    "type": "string",
                                    "description": "Document title",
                                },
                                "content": {
                                    "type": "string",
                                    "description": "Document content (supports markdown-like formatting)",
                                },
                            },
                            "required": ["file_path", "content"],
                        },
                        handler=self.word_create,
                    ),
                    Tool(
                        name="word_add_paragraph",
                        description="Add a paragraph to an existing Word document",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .docx file",
                                },
                                "text": {
                                    "type": "string",
                                    "description": "Text to add",
                                },
                                "style": {
                                    "type": "string",
                                    "description": "Style: Normal, Heading 1, Heading 2, Title, Quote",
                                },
                            },
                            "required": ["file_path", "text"],
                        },
                        handler=self.word_add_paragraph,
                    ),
                    Tool(
                        name="word_add_table",
                        description="Add a table to a Word document",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .docx file",
                                },
                                "data": {
                                    "type": "array",
                                    "description": "Table data as array of arrays [[row1], [row2], ...]",
                                    "items": {"type": "array"},
                                },
                                "headers": {
                                    "type": "boolean",
                                    "description": "First row is headers (default: true)",
                                },
                            },
                            "required": ["file_path", "data"],
                        },
                        handler=self.word_add_table,
                    ),
                ]
            )

        # Excel tools
        if HAS_XLSX:
            tools.extend(
                [
                    Tool(
                        name="excel_read",
                        description="Read content from an Excel spreadsheet (.xlsx)",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .xlsx file",
                                },
                                "sheet_name": {
                                    "type": "string",
                                    "description": "Sheet name (default: active sheet)",
                                },
                                "range": {
                                    "type": "string",
                                    "description": "Cell range like 'A1:D10' (default: all data)",
                                },
                            },
                            "required": ["file_path"],
                        },
                        handler=self.excel_read,
                    ),
                    Tool(
                        name="excel_create",
                        description="Create a new Excel spreadsheet with data",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Output path for the .xlsx file",
                                },
                                "data": {
                                    "type": "array",
                                    "description": "Data as array of arrays [[row1], [row2], ...]",
                                    "items": {"type": "array"},
                                },
                                "sheet_name": {
                                    "type": "string",
                                    "description": "Sheet name (default: Sheet1)",
                                },
                                "headers": {
                                    "type": "boolean",
                                    "description": "Style first row as headers (default: true)",
                                },
                            },
                            "required": ["file_path", "data"],
                        },
                        handler=self.excel_create,
                    ),
                    Tool(
                        name="excel_write_cell",
                        description="Write value to a specific cell in Excel",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .xlsx file",
                                },
                                "cell": {
                                    "type": "string",
                                    "description": "Cell reference like 'A1' or 'B5'",
                                },
                                "value": {
                                    "description": "Value to write",
                                },
                                "sheet_name": {
                                    "type": "string",
                                    "description": "Sheet name (default: active sheet)",
                                },
                            },
                            "required": ["file_path", "cell", "value"],
                        },
                        handler=self.excel_write_cell,
                    ),
                    Tool(
                        name="excel_add_sheet",
                        description="Add a new sheet to an Excel workbook",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .xlsx file",
                                },
                                "sheet_name": {
                                    "type": "string",
                                    "description": "Name for the new sheet",
                                },
                                "data": {
                                    "type": "array",
                                    "description": "Optional data to populate",
                                    "items": {"type": "array"},
                                },
                            },
                            "required": ["file_path", "sheet_name"],
                        },
                        handler=self.excel_add_sheet,
                    ),
                    Tool(
                        name="excel_list_sheets",
                        description="List all sheets in an Excel workbook",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .xlsx file",
                                },
                            },
                            "required": ["file_path"],
                        },
                        handler=self.excel_list_sheets,
                    ),
                    Tool(
                        name="excel_formula",
                        description="Add a formula to a cell",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .xlsx file",
                                },
                                "cell": {
                                    "type": "string",
                                    "description": "Cell reference like 'A1'",
                                },
                                "formula": {
                                    "type": "string",
                                    "description": "Excel formula like '=SUM(A1:A10)'",
                                },
                                "sheet_name": {
                                    "type": "string",
                                    "description": "Sheet name (default: active sheet)",
                                },
                            },
                            "required": ["file_path", "cell", "formula"],
                        },
                        handler=self.excel_formula,
                    ),
                ]
            )

        # PowerPoint tools
        if HAS_PPTX:
            tools.extend(
                [
                    Tool(
                        name="pptx_read",
                        description="Read content from a PowerPoint presentation (.pptx)",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .pptx file",
                                },
                            },
                            "required": ["file_path"],
                        },
                        handler=self.pptx_read,
                    ),
                    Tool(
                        name="pptx_create",
                        description="Create a new PowerPoint presentation",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Output path for the .pptx file",
                                },
                                "title": {
                                    "type": "string",
                                    "description": "Presentation title",
                                },
                                "subtitle": {
                                    "type": "string",
                                    "description": "Presentation subtitle",
                                },
                            },
                            "required": ["file_path", "title"],
                        },
                        handler=self.pptx_create,
                    ),
                    Tool(
                        name="pptx_add_slide",
                        description="Add a slide to a PowerPoint presentation",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .pptx file",
                                },
                                "title": {
                                    "type": "string",
                                    "description": "Slide title",
                                },
                                "content": {
                                    "type": "string",
                                    "description": "Slide content (bullet points separated by newlines)",
                                },
                                "layout": {
                                    "type": "string",
                                    "description": "Layout: title, content, two_content, blank",
                                },
                            },
                            "required": ["file_path", "title"],
                        },
                        handler=self.pptx_add_slide,
                    ),
                    Tool(
                        name="pptx_add_image",
                        description="Add an image to a PowerPoint slide",
                        parameters={
                            "type": "object",
                            "properties": {
                                "file_path": {
                                    "type": "string",
                                    "description": "Path to the .pptx file",
                                },
                                "image_path": {
                                    "type": "string",
                                    "description": "Path to the image file",
                                },
                                "slide_index": {
                                    "type": "integer",
                                    "description": "Slide index (0-based, default: last slide)",
                                },
                            },
                            "required": ["file_path", "image_path"],
                        },
                        handler=self.pptx_add_image,
                    ),
                ]
            )

        return tools

    # ==================== WORD METHODS ====================

    def word_read(self, file_path: str, include_tables: bool = True) -> str:
        """Read content from a Word document."""
        if not HAS_DOCX:
            return "Error: python-docx not installed. Run: pip install python-docx"

        try:
            path = Path(file_path).expanduser()
            if not path.exists():
                return f"File not found: {file_path}"

            doc = Document(path)
            content = []

            # Read paragraphs
            for para in doc.paragraphs:
                if para.text.strip():
                    style = para.style.name if para.style else "Normal"
                    if "Heading" in style:
                        content.append(f"\n## {para.text}\n")
                    elif style == "Title":
                        content.append(f"\n# {para.text}\n")
                    else:
                        content.append(para.text)

            # Read tables
            if include_tables and doc.tables:
                for i, table in enumerate(doc.tables):
                    content.append(f"\n[Table {i + 1}]")
                    for row in table.rows:
                        cells = [cell.text.strip() for cell in row.cells]
                        content.append(" | ".join(cells))

            result = f"=== Word Document: {path.name} ===\n\n"
            result += "\n".join(content)
            return result

        except Exception as e:
            return f"Error reading Word document: {e}"

    def word_create(
        self,
        file_path: str,
        content: str,
        title: str | None = None,
    ) -> str:
        """Create a new Word document."""
        if not HAS_DOCX:
            return "Error: python-docx not installed. Run: pip install python-docx"

        try:
            path = Path(file_path).expanduser()
            path.parent.mkdir(parents=True, exist_ok=True)

            doc = Document()

            # Add title if provided
            if title:
                doc.add_heading(title, 0)

            # Parse content (simple markdown-like)
            lines = content.split("\n")
            for line in lines:
                line = line.strip()
                if not line:
                    continue

                if line.startswith("# "):
                    doc.add_heading(line[2:], 1)
                elif line.startswith("## "):
                    doc.add_heading(line[3:], 2)
                elif line.startswith("### "):
                    doc.add_heading(line[4:], 3)
                elif line.startswith("- ") or line.startswith("* "):
                    doc.add_paragraph(line[2:], style="List Bullet")
                elif line.startswith("1. ") or (line[0].isdigit() and line[1] == "."):
                    doc.add_paragraph(line[3:], style="List Number")
                elif line.startswith("> "):
                    para = doc.add_paragraph(line[2:])
                    para.style = "Quote"
                else:
                    doc.add_paragraph(line)

            doc.save(path)
            return f"Word document created: {path}"

        except Exception as e:
            return f"Error creating Word document: {e}"

    def word_add_paragraph(
        self,
        file_path: str,
        text: str,
        style: str = "Normal",
    ) -> str:
        """Add a paragraph to a Word document."""
        if not HAS_DOCX:
            return "Error: python-docx not installed. Run: pip install python-docx"

        try:
            path = Path(file_path).expanduser()
            if not path.exists():
                return f"File not found: {file_path}"

            doc = Document(path)
            doc.add_paragraph(text, style=style)
            doc.save(path)

            return f"Paragraph added to {path.name}"

        except Exception as e:
            return f"Error adding paragraph: {e}"

    def word_add_table(
        self,
        file_path: str,
        data: list[list],
        headers: bool = True,
    ) -> str:
        """Add a table to a Word document."""
        if not HAS_DOCX:
            return "Error: python-docx not installed. Run: pip install python-docx"

        try:
            path = Path(file_path).expanduser()
            if not path.exists():
                return f"File not found: {file_path}"

            doc = Document(path)

            if not data:
                return "Error: No data provided for table"

            rows = len(data)
            cols = len(data[0]) if data else 0

            table = doc.add_table(rows=rows, cols=cols)
            table.style = "Table Grid"

            for i, row_data in enumerate(data):
                row = table.rows[i]
                for j, cell_data in enumerate(row_data):
                    cell = row.cells[j]
                    cell.text = str(cell_data)

                    # Bold headers
                    if headers and i == 0:
                        for paragraph in cell.paragraphs:
                            for run in paragraph.runs:
                                run.bold = True

            doc.save(path)
            return f"Table ({rows}x{cols}) added to {path.name}"

        except Exception as e:
            return f"Error adding table: {e}"

    # ==================== EXCEL METHODS ====================

    def excel_read(
        self,
        file_path: str,
        sheet_name: str | None = None,
        range: str | None = None,
    ) -> str:
        """Read content from an Excel file."""
        if not HAS_XLSX:
            return "Error: openpyxl not installed. Run: pip install openpyxl"

        try:
            path = Path(file_path).expanduser()
            if not path.exists():
                return f"File not found: {file_path}"

            wb = load_workbook(path, data_only=True)
            ws = wb[sheet_name] if sheet_name else wb.active

            result = f"=== Excel: {path.name} ({ws.title}) ===\n\n"

            if range:
                # Read specific range
                cells = ws[range]
                if hasattr(cells, "__iter__"):
                    for row in cells:
                        if hasattr(row, "__iter__"):
                            values = [str(cell.value) if cell.value else "" for cell in row]
                            result += " | ".join(values) + "\n"
                        else:
                            result += str(row.value) + "\n"
            else:
                # Read all data
                for row in ws.iter_rows(values_only=True):
                    if any(cell is not None for cell in row):
                        values = [str(cell) if cell is not None else "" for cell in row]
                        result += " | ".join(values) + "\n"

            wb.close()
            return result

        except Exception as e:
            return f"Error reading Excel file: {e}"

    def excel_create(
        self,
        file_path: str,
        data: list[list],
        sheet_name: str = "Sheet1",
        headers: bool = True,
    ) -> str:
        """Create a new Excel file."""
        if not HAS_XLSX:
            return "Error: openpyxl not installed. Run: pip install openpyxl"

        try:
            path = Path(file_path).expanduser()
            path.parent.mkdir(parents=True, exist_ok=True)

            wb = Workbook()
            ws = wb.active
            ws.title = sheet_name

            # Header style
            header_font = Font(bold=True)
            header_fill = PatternFill(start_color="DDEEFF", end_color="DDEEFF", fill_type="solid")

            for row_idx, row_data in enumerate(data, 1):
                for col_idx, value in enumerate(row_data, 1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=value)

                    # Style headers
                    if headers and row_idx == 1:
                        cell.font = header_font
                        cell.fill = header_fill

            # Auto-adjust column widths
            for col in ws.columns:
                max_length = 0
                column = col[0].column_letter
                for cell in col:
                    if cell.value:
                        max_length = max(max_length, len(str(cell.value)))
                ws.column_dimensions[column].width = min(max_length + 2, 50)

            wb.save(path)
            wb.close()

            rows = len(data)
            cols = len(data[0]) if data else 0
            return f"Excel file created: {path} ({rows} rows, {cols} columns)"

        except Exception as e:
            return f"Error creating Excel file: {e}"

    def excel_write_cell(
        self,
        file_path: str,
        cell: str,
        value,
        sheet_name: str | None = None,
    ) -> str:
        """Write a value to a specific cell."""
        if not HAS_XLSX:
            return "Error: openpyxl not installed. Run: pip install openpyxl"

        try:
            path = Path(file_path).expanduser()

            if path.exists():
                wb = load_workbook(path)
            else:
                wb = Workbook()

            ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
            ws[cell] = value
            wb.save(path)
            wb.close()

            return f"Cell {cell} updated to: {value}"

        except Exception as e:
            return f"Error writing cell: {e}"

    def excel_add_sheet(
        self,
        file_path: str,
        sheet_name: str,
        data: list[list] | None = None,
    ) -> str:
        """Add a new sheet to an Excel file."""
        if not HAS_XLSX:
            return "Error: openpyxl not installed. Run: pip install openpyxl"

        try:
            path = Path(file_path).expanduser()

            if path.exists():
                wb = load_workbook(path)
            else:
                wb = Workbook()

            ws = wb.create_sheet(title=sheet_name)

            if data:
                for row_idx, row_data in enumerate(data, 1):
                    for col_idx, value in enumerate(row_data, 1):
                        ws.cell(row=row_idx, column=col_idx, value=value)

            wb.save(path)
            wb.close()

            return f"Sheet '{sheet_name}' added to {path.name}"

        except Exception as e:
            return f"Error adding sheet: {e}"

    def excel_list_sheets(self, file_path: str) -> str:
        """List all sheets in an Excel file."""
        if not HAS_XLSX:
            return "Error: openpyxl not installed. Run: pip install openpyxl"

        try:
            path = Path(file_path).expanduser()
            if not path.exists():
                return f"File not found: {file_path}"

            wb = load_workbook(path)
            sheets = wb.sheetnames
            wb.close()

            result = f"Sheets in {path.name}:\n"
            for i, sheet in enumerate(sheets, 1):
                result += f"  {i}. {sheet}\n"

            return result

        except Exception as e:
            return f"Error listing sheets: {e}"

    def excel_formula(
        self,
        file_path: str,
        cell: str,
        formula: str,
        sheet_name: str | None = None,
    ) -> str:
        """Add a formula to a cell."""
        if not HAS_XLSX:
            return "Error: openpyxl not installed. Run: pip install openpyxl"

        try:
            path = Path(file_path).expanduser()

            if path.exists():
                wb = load_workbook(path)
            else:
                wb = Workbook()

            ws = wb[sheet_name] if sheet_name and sheet_name in wb.sheetnames else wb.active
            ws[cell] = formula
            wb.save(path)
            wb.close()

            return f"Formula added to {cell}: {formula}"

        except Exception as e:
            return f"Error adding formula: {e}"

    # ==================== POWERPOINT METHODS ====================

    def pptx_read(self, file_path: str) -> str:
        """Read content from a PowerPoint file."""
        if not HAS_PPTX:
            return "Error: python-pptx not installed. Run: pip install python-pptx"

        try:
            path = Path(file_path).expanduser()
            if not path.exists():
                return f"File not found: {file_path}"

            prs = Presentation(path)
            result = f"=== PowerPoint: {path.name} ===\n"
            result += f"Slides: {len(prs.slides)}\n\n"

            for i, slide in enumerate(prs.slides, 1):
                result += f"--- Slide {i} ---\n"

                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        result += f"{shape.text}\n"

                result += "\n"

            return result

        except Exception as e:
            return f"Error reading PowerPoint: {e}"

    def pptx_create(
        self,
        file_path: str,
        title: str,
        subtitle: str | None = None,
    ) -> str:
        """Create a new PowerPoint presentation."""
        if not HAS_PPTX:
            return "Error: python-pptx not installed. Run: pip install python-pptx"

        try:
            path = Path(file_path).expanduser()
            path.parent.mkdir(parents=True, exist_ok=True)

            prs = Presentation()

            # Title slide
            title_slide_layout = prs.slide_layouts[0]
            slide = prs.slides.add_slide(title_slide_layout)

            title_shape = slide.shapes.title
            title_shape.text = title

            if subtitle:
                subtitle_shape = slide.placeholders[1]
                subtitle_shape.text = subtitle

            prs.save(path)
            return f"PowerPoint created: {path}"

        except Exception as e:
            return f"Error creating PowerPoint: {e}"

    def pptx_add_slide(
        self,
        file_path: str,
        title: str,
        content: str | None = None,
        layout: str = "content",
    ) -> str:
        """Add a slide to a PowerPoint presentation."""
        if not HAS_PPTX:
            return "Error: python-pptx not installed. Run: pip install python-pptx"

        try:
            path = Path(file_path).expanduser()

            if path.exists():
                prs = Presentation(path)
            else:
                prs = Presentation()

            # Layout mapping
            layout_map = {
                "title": 0,
                "content": 1,
                "two_content": 3,
                "blank": 6,
            }
            layout_idx = layout_map.get(layout, 1)

            slide_layout = prs.slide_layouts[layout_idx]
            slide = prs.slides.add_slide(slide_layout)

            # Set title
            if slide.shapes.title:
                slide.shapes.title.text = title

            # Set content
            if content and len(slide.placeholders) > 1:
                body_shape = slide.placeholders[1]
                tf = body_shape.text_frame

                lines = content.split("\n")
                for i, line in enumerate(lines):
                    if i == 0:
                        tf.text = line.strip()
                    else:
                        p = tf.add_paragraph()
                        p.text = line.strip()
                        p.level = 0

            prs.save(path)
            slide_num = len(prs.slides)
            return f"Slide {slide_num} added to {path.name}: {title}"

        except Exception as e:
            return f"Error adding slide: {e}"

    def pptx_add_image(
        self,
        file_path: str,
        image_path: str,
        slide_index: int = -1,
    ) -> str:
        """Add an image to a slide."""
        if not HAS_PPTX:
            return "Error: python-pptx not installed. Run: pip install python-pptx"

        try:
            pptx_path = Path(file_path).expanduser()
            img_path = Path(image_path).expanduser()

            if not pptx_path.exists():
                return f"PowerPoint not found: {file_path}"
            if not img_path.exists():
                return f"Image not found: {image_path}"

            prs = Presentation(pptx_path)

            if slide_index < 0:
                slide_index = len(prs.slides) - 1

            if slide_index >= len(prs.slides):
                return f"Slide index {slide_index} out of range"

            slide = prs.slides[slide_index]

            # Add image (centered, 5 inches wide)
            left = PptxInches(1)
            top = PptxInches(2)
            width = PptxInches(5)

            slide.shapes.add_picture(str(img_path), left, top, width=width)

            prs.save(pptx_path)
            return f"Image added to slide {slide_index + 1}"

        except Exception as e:
            return f"Error adding image: {e}"

    def execute(self, **kwargs) -> str:
        """Direct execution."""
        action = kwargs.get("action", "info")

        if action == "info":
            result = "=== MS Office Skill ===\n\n"
            result += f"Word (.docx): {'Available' if HAS_DOCX else 'Not installed'}\n"
            result += f"Excel (.xlsx): {'Available' if HAS_XLSX else 'Not installed'}\n"
            result += f"PowerPoint (.pptx): {'Available' if HAS_PPTX else 'Not installed'}\n"
            return result
        else:
            return f"Unknown action: {action}"
